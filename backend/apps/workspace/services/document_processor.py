import hashlib
import logging
import os
import tempfile

from apps.workspace.models import Document, DocumentChunk
from apps.workspace.services.chunk_service import ChunkService
from apps.workspace.services.embedding_service import EmbeddingService
from apps.workspace.services.chromadb_service import VectorStoreService
from apps.workspace.services.supabase_service import SupabaseStorageService
from apps.workspace.exceptions import RetryableProcessingError, NonRetryableProcessingError

logger = logging.getLogger(__name__)


class DocumentProcessor:

    @staticmethod
    def extract_text(document):
        extension = document.file.name.split('.')[-1].lower()
        file_bytes = SupabaseStorageService.download_file(document.file.name)
        temp_path = None

        try:
            with tempfile.NamedTemporaryFile(
                delete=False,
                suffix=f".{extension}"
            ) as temp_file:
                temp_file.write(file_bytes)
                temp_path = temp_file.name

            if extension in ['txt', 'html', 'md']:
                with open(temp_path, 'r', encoding='utf-8') as file:
                    return file.read()

            elif extension == 'pdf':
                from pypdf import PdfReader

                reader = PdfReader(temp_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() or ""
                return text

            elif extension == 'docx':
                from docx import Document as DocxDocument

                doc = DocxDocument(temp_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text

            elif extension == 'xlsx':
                from openpyxl import load_workbook

                workbook = load_workbook(temp_path, data_only=True)
                text = ""
                for sheet in workbook.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join(str(cell) for cell in row if cell is not None)
                        text += row_text + "\n"
                return text

            elif extension == 'pptx':
                from pptx import Presentation

                presentation = Presentation(temp_path)
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text

            elif extension == 'csv':
                import csv

                text = ""
                with open(temp_path, 'r', encoding='utf-8') as file:
                    reader = csv.reader(file)
                    for row in reader:
                        text += " ".join(row) + "\n"
                return text

            return None

        finally:
            if temp_path and os.path.exists(temp_path):
                os.remove(temp_path)


    @staticmethod
    def update_status(document_id, status):
        document = Document.objects.get(id=document_id)
        document.status = status
        document.save(update_fields=["status"])


    @staticmethod
    def process(document_id):
        document = None

        try:
            document = Document.objects.get(id=document_id)
            text = DocumentProcessor.extract_text(document)

            if not text:
                DocumentProcessor.update_status(document_id, "not_supported")
                return

            content_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()
            existing_document = Document.objects.filter(
                user=document.user,
                content_hash=content_hash
            ).exclude(id=document.id).exists()

            if existing_document:
                document.delete()
                logger.warning(f"Document Already Exists: {document_id}")
                return False

            logger.info(f"Started processing document: {document_id}")

            document.content_hash = content_hash
            document.extracted_data = text
            document.save()

            chunks = ChunkService.create_chunks(text)

            chunk_objects = [
                DocumentChunk(
                    document=document,
                    chunk_text=chunk,
                    chunk_id=index
                )
                for index, chunk in enumerate(chunks)
            ]

            DocumentChunk.objects.bulk_create(chunk_objects)
            logger.info("Chunk generated Successfully")

            embeddings = EmbeddingService.generate_embeddings(chunks)
            logger.info("Embeddings generated successfully.")

            VectorStoreService.add_chunks(document, chunks, embeddings)
            logger.info("Vectors generated successfully")

            DocumentProcessor.update_status(document_id, "ready")
            logger.info("Document processing completed successfully.")

            return True

        except Document.DoesNotExist as e:
            logger.error(f"Document {document_id} not found.")
            raise NonRetryableProcessingError("Document does not exist.") from e

        except FileNotFoundError as e:
            logger.error(f"Document file not found: {document_id}")
            raise NonRetryableProcessingError("Document file not found.") from e

        except PermissionError as e:
            logger.error(f"Permission denied while reading document {document_id}")
            raise NonRetryableProcessingError("Permission denied.") from e

        except Exception as e:
            logger.exception(f"Unexpected error while processing document {document_id}")
            raise RetryableProcessingError("Temporary failure during document processing.") from e


    @staticmethod
    def cleanup_artifacts(document_id):
        document = Document.objects.get(id=document_id)
        DocumentChunk.objects.filter(document=document).delete()
        VectorStoreService.delete_vector(document=document)
        logger.info(f"Previous artifacts removed successfully for {document_id}")